# -*- coding: utf-8 -*-
"""Generador PPTX — Recepción de Vehículos · Planta Entrerrios
v2.0 — Dashboard pastel, etiquetas visibles (DPI 220), logos Alpina embebidos,
navegación interactiva entre slides (hipervínculos internos).
"""
import re, os, io, base64, tempfile
from datetime import date
import matplotlib
matplotlib.use("Agg")

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ─── Paleta PASTEL Operativa ────────────────────────────────────────────────
OP_AZ  = RGBColor(0x1D,0x4E,0xD8)
OP_AZ2 = RGBColor(0x60,0xA5,0xFA)
OP_AZL = RGBColor(0xDB,0xEA,0xFE)
OP_BG  = RGBColor(0xF0,0xF7,0xFF)
OP_VD  = RGBColor(0x16,0xA3,0x4A)
OP_VDL = RGBColor(0xD1,0xFA,0xE5)
OP_RJ  = RGBColor(0xDC,0x26,0x26)
OP_RJL = RGBColor(0xFE,0xE2,0xE2)
OP_NR  = RGBColor(0xEA,0x58,0x0C)
OP_NRL = RGBColor(0xFF,0xED,0xD5)
OP_WH  = RGBColor(0xFF,0xFF,0xFF)
OP_GR  = RGBColor(0x64,0x74,0x8B)
OP_DARK= RGBColor(0x0F,0x1C,0x3A)

# ─── Paleta PASTEL Gerencial ────────────────────────────────────────────────
GR_BG  = RGBColor(0xFF,0xFB,0xF0)
GR_GD  = RGBColor(0x92,0x40,0x0E)
GR_GDL = RGBColor(0xFE,0xF3,0xC7)
GR_AZ  = RGBColor(0x1E,0x3A,0x8A)
GR_AZL = RGBColor(0xDB,0xEA,0xFE)
GR_VD  = RGBColor(0x16,0xA3,0x4A)
GR_RJ  = RGBColor(0xDC,0x26,0x26)
GR_WH  = RGBColor(0xFF,0xFF,0xFF)

# ─── Matplotlib colores pastel ───────────────────────────────────────────────
PASTEL = {
    "azul"   : "#90CAF9", "verde"  : "#A5D6A7", "naranja": "#FFCC80",
    "rojo"   : "#EF9A9A", "morado" : "#CE93D8", "celeste": "#80DEEA",
    "amarillo":"#FFF176", "gris"   : "#CFD8DC"
}

# ─── Imágenes Base64 (reemplaza con tus cadenas completas) ──────────────────
_IMG2_B64 = "" # ← PEGA AQUÍ TU BASE64 ORIGINAL (Logo azul)
_IMG4_B64 = "" # ← PEGA AQUÍ TU BASE64 ORIGINAL (Logo esquina)
_TRUCK_B64= "" # ← PEGA AQUÍ TU BASE64 ORIGINAL (Camión)

## ─── Helpers PPTX ────────────────────────────────────────────────────────────

def _img_place_file(slide, ruta_img, x, y, w, h):
    if not ruta_img or not os.path.exists(ruta_img):
        print(f"Imagen no encontrada: {ruta_img}")
        return

    try:
        slide.shapes.add_picture(
            ruta_img,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h)
        )
    except Exception as e:
        print(f"Error cargando imagen {ruta_img}: {e}")


def _img_place_stream(slide, stream, x, y, w, h):
    try:
        slide.shapes.add_picture(
            stream,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h)
        )
    except Exception as e:
        print(f"Error cargando imagen desde stream: {e}")


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, color, radius=0):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _txt(slide, text, x, y, w, h, size=12, bold=False,
         color=RGBColor(255, 255, 255),
         align=PP_ALIGN.LEFT,
         italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap

    p = tf.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    # Filtrar caracteres fuera del BMP (>U+FFFF) que pueden causar errores XML en python-pptx
    safe_text = "".join(c if ord(c) <= 0xFFFF else "?" for c in str(text))
    run.text = safe_text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"

    return txb


def _add_logo(slide):
    try:
        _img_place_file(slide, "assets/Imagen4.jpg", 11.9, 0.08, 1.28, 0.62)
    except Exception:
        pass


def _add_nav(slide, prs, idx):
    total = len(prs.slides)

    if idx > 0:
        shp = slide.shapes.add_shape(
            1,
            Inches(0.25),
            Inches(7.1),
            Inches(0.55),
            Inches(0.28)
        )

        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(59, 130, 246)
        shp.line.fill.background()

        tb = shp.text_frame
        tb.word_wrap = False

        p = tb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        r = p.add_run()
        r.text = "<"
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(255, 255, 255)
        r.font.name = "Calibri"

        _add_slide_hyperlink(shp, idx - 1)


def _add_slide_hyperlink(shape, target_idx):
    try:
        sp = shape._element
        # Para un auto-shape, cNvPr está en nvSpPr/cNvPr
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            return
        cPr = nvSpPr.find(qn("p:cNvPr"))
        if cPr is None:
            return

        hlinkClick = etree.SubElement(
            cPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick"
        )
        hlinkClick.attrib["action"] = f"ppaction://hlinksldjump?slideindex={target_idx}&return=false"

    except Exception:
        pass

# ─── Helpers matplotlib ──────────────────────────────────────────────────────
def _ax_style(ax, bg="#F0F7FF", gc="#DDE6F0"):
    ax.set_facecolor(bg)
    ax.spines[["top","right","left","bottom"]].set_visible(False)
    ax.tick_params(colors="#374151", labelsize=8.5, length=0)
    ax.yaxis.grid(True, color=gc, linewidth=0.6, linestyle="--")
    ax.set_axisbelow(True)

def _ax_style_warm(ax):
    _ax_style(ax, "#FFF8F0", "#F0E0D0")

def _fig_to_stream(fig, dpi=220):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor(), dpi=dpi)
    buf.seek(0)
    plt.close(fig)
    return buf

def _bar_labels(ax, fmt="{:.0f}", color="#374151", fontsize=8.5, pad=1.5):
    for container in ax.containers:
        ax.bar_label(container, fmt=fmt, padding=pad, fontsize=fontsize, color=color, fontweight="bold")

def _barh_labels(ax, vals, fmt="{:.0f}", color="#374151", fontsize=8, pad=1.5):
    for i, v in enumerate(vals):
        ax.text(v + pad, i, fmt.format(v), va="center", ha="left", fontsize=fontsize, color=color, fontweight="bold")

def _kpi_box(slide, val, lbl, sub, x, y, bg=OP_AZ2, val_c=OP_WH, lbl_c=OP_GR, w=2.0, h=1.45):
    _rect(slide, x, y, w, h, bg)
    _rect(slide, x, y, w, 0.07, val_c)
    _txt(slide, str(val), x+0.07, y+0.15, w-0.14, 0.68, size=26, bold=True, color=val_c, align=PP_ALIGN.CENTER)
    _txt(slide, lbl, x+0.05, y+0.76, w-0.1, 0.36, size=9, bold=True, color=lbl_c, align=PP_ALIGN.CENTER)
    if sub:
        _txt(slide, sub, x+0.05, y+1.09, w-0.1, 0.28, size=7.5, color=OP_GR, align=PP_ALIGN.CENTER, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# PRESENTACIÓN OPERATIVA
# ═════════════════════════════════════════════════════════════════════════════
def generar_pptx_operativo(df_ag, df_cc, label, fd, fh, carpeta):
    prs = _new_prs()
    total = len(df_ag)
    recib = int((df_ag["_estado"]=="RECIBIDO").sum())
    rechaz = int((df_ag["_estado"]=="RECHAZADO").sum())
    no_ll = int(df_ag["_estado"].str.contains("LLEG",na=False).sum())
    pct_ef = recib/total if total else 0
    est_tot = int(df_ag["_estibas"].sum())
    min_ok = df_cc[df_cc["_min"]>0]["_min"]
    prom_t = min_ok.mean() if len(min_ok) else 0
    fecha_g = date.today().strftime("%d/%m/%Y")
    sin_placa = int(df_cc["_novedad"].str.strip().eq("SIN PLACA").sum())
    no_msd = int(df_cc["_novedad"].str.strip().eq("NO DEJA AGENDAR MSD").sum())
    llega_t = int(df_cc["_novedad"].str.strip().eq("LLEGA TARDE").sum())
    slides_list = []

    # SLIDE 1: PORTADA
    sl = _blank(prs); slides_list.append(sl)
    _rect(sl, 0, 0, 13.33, 7.5, OP_BG)
    _rect(sl, 0, 0, 13.33, 3.4, OP_AZ)
    _rect(sl, 0, 0, 0.22, 7.5, OP_AZ2)
    _img_place_file(sl, "assets/Imagen1.png", 0.5, 0.3, 1.85, 2.6)
    _txt(sl, "INFORME OPERATIVO", 2.6, 0.4, 10.3, 0.9, size=33, bold=True, color=OP_WH, align=PP_ALIGN.LEFT)
    _txt(sl, "RECEPCION DE VEHICULOS - PLANTA ENTRERRIOS", 2.6, 1.32, 10.3, 0.48, size=15, color=RGBColor(0xBA,0xD8,0xF8), align=PP_ALIGN.LEFT)
    _txt(sl, f"Periodo: {label}", 2.6, 1.85, 10.3, 0.42, size=13, color=OP_AZL, align=PP_ALIGN.LEFT)
    _txt(sl, f"{fd.strftime('%d/%m/%Y')} - {fh.strftime('%d/%m/%Y')}", 2.6, 2.3, 10.3, 0.38, size=11, color=OP_GR, align=PP_ALIGN.LEFT)
    _rect(sl, 0.35, 3.55, 12.63, 0.05, OP_AZ2)
    idx_items = ["1 - Portada","2 - Resumen KPIs","3 - Evolucion Diaria","4 - Analisis Proveedores","5 - Novedades & Tiempos"]
    for i, item in enumerate(idx_items):
        yi = 3.75 + i * 0.58
        _rect(sl, 0.5, yi, 0.05, 0.38, OP_AZ2)
        _txt(sl, item, 0.72, yi+0.03, 5.5, 0.35, size=11, color=OP_AZ, align=PP_ALIGN.LEFT)
    _txt(sl, f"Generado: {fecha_g} | Agenda: {total} reg. | Control Citas: {len(df_cc)} reg.", 0.35, 7.1, 12.63, 0.3, size=8.5, color=OP_GR, align=PP_ALIGN.CENTER, italic=True)
    _add_logo(sl); _add_nav(sl, prs, 0)

    # SLIDE 2: KPI RESUMEN
    sl = _blank(prs); slides_list.append(sl)
    _rect(sl, 0, 0, 13.33, 7.5, OP_BG)
    _rect(sl, 0, 0, 0.22, 7.5, OP_AZ2)
    _rect(sl, 0.22, 0, 13.11, 0.62, OP_AZ)
    _txt(sl, " RESUMEN DE METRICAS - PERIODO", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=OP_WH, align=PP_ALIGN.LEFT)
    _txt(sl, f"{label} · {fd.strftime('%d/%m/%Y')} al {fh.strftime('%d/%m/%Y')}", 0.35, 0.65, 12, 0.3, size=9.5, color=OP_GR, align=PP_ALIGN.CENTER, italic=True)
    _img_place_file(sl, "assets/Gemini_Generated_Image_vcna55vcna55vcna.png", 7.8, 0.95, 5.4, 5.8)
    ef_col = OP_VD if pct_ef>=.85 else (RGBColor(0xEA,0x58,0x0C) if pct_ef>=.70 else OP_RJ)
    ef_bg = RGBColor(0xD1,0xFA,0xE5) if pct_ef>=.85 else (RGBColor(0xFF,0xED,0xD5) if pct_ef>=.70 else RGBColor(0xFE,0xE2,0xE2))
    kpis = [
        (f"{pct_ef:.1%}", "Efectividad", "Meta >=85%", ef_bg, ef_col),
        (str(total), "Agendados", "Total periodo", OP_AZL, OP_AZ),
        (str(recib), "Recibidos", f"{recib/total:.0%} del total" if total else "-", RGBColor(0xD1,0xFA,0xE5), OP_VD),
        (str(no_ll), "No Llego", f"{no_ll/total:.0%} del total" if total else "-", RGBColor(0xFF,0xED,0xD5), OP_NR),
        (str(rechaz), "Rechazados", f"{rechaz/total:.0%}" if total else "-", RGBColor(0xFE,0xE2,0xE2), OP_RJ),
        (f"{prom_t:.0f} min", "Prom. Descargue", "Limite 60 min", RGBColor(0xF0,0xF4,0xFF), OP_AZ),
        (f"{est_tot:,}", "Estibas", "Acumuladas", OP_AZL, OP_AZ),
        (str(len(df_cc)), "Reg. Control CC", "Registros cargados", OP_AZL, OP_AZ),
    ]
    cols_per_row = 4
    for i, (v, l, s, bg, fc) in enumerate(kpis):
        row, col_i = divmod(i, cols_per_row)
        xi = 0.4 + col_i * 1.85
        yi = 1.05 + row * 1.62
        _rect(sl, xi, yi, 1.72, 1.48, bg)
        _rect(sl, xi, yi, 1.72, 0.07, fc)
        _txt(sl, v, xi+0.06, yi+0.12, 1.6, 0.68, size=26, bold=True, color=fc, align=PP_ALIGN.CENTER)
        _txt(sl, l, xi+0.04, yi+0.76, 1.64, 0.34, size=8.5, bold=True, color=OP_GR, align=PP_ALIGN.CENTER)
        if s:
            _txt(sl, s, xi+0.04, yi+1.07, 1.64, 0.3, size=7, color=RGBColor(0x9C,0xA3,0xAF), align=PP_ALIGN.CENTER, italic=True)
    _txt(sl, "* Arrastra el raton o usa los botones para navegar", 0.35, 7.1, 12.63, 0.3, size=7.5, color=RGBColor(0x9C,0xA3,0xAF), align=PP_ALIGN.CENTER, italic=True)
    _add_logo(sl); _add_nav(sl, prs, 1)

    # SLIDE 3: EVOLUCIÓN DIARIA
    sl = _blank(prs); slides_list.append(sl)
    _rect(sl, 0, 0, 13.33, 7.5, OP_BG)
    _rect(sl, 0, 0, 0.22, 7.5, OP_AZ2)
    _rect(sl, 0.22, 0, 13.11, 0.62, OP_AZ)
    _txt(sl, " EVOLUCION DIARIA DE RECEPCIONES", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=OP_WH, align=PP_ALIGN.LEFT)
    _txt(sl, f"{label} · {fd.strftime('%d/%m/%Y')} al {fh.strftime('%d/%m/%Y')}", 0.35, 0.65, 12, 0.3, size=9.5, color=OP_GR, align=PP_ALIGN.CENTER, italic=True)
    df_d = df_ag.copy(); df_d["_ds"] = df_d["_fecha"].astype(str)
    ev = df_d.groupby("_ds").agg(
        total=("_estado","count"),
        recibidos=("_estado", lambda x:(x=="RECIBIDO").sum()),
        no_llego=("_estado", lambda x:x.str.contains("LLEG",na=False).sum()),
        rechazados=("_estado",lambda x:(x=="RECHAZADO").sum())
    ).reset_index().sort_values("_ds")
    xs = np.arange(len(ev))
    lbl = [d[5:] for d in ev["_ds"]]
    pct_d = [r/t*100 if t else 0 for r,t in zip(ev["recibidos"],ev["total"])]
    fig, (ax1,ax2) = plt.subplots(2, 1, figsize=(12.5, 5.8), facecolor="#F0F7FF")
    fig.subplots_adjust(hspace=0.55)
    width = 0.35
    b1 = ax1.bar(xs-width/2, ev["total"], width, label="Agendados", color=PASTEL["azul"], alpha=0.9)
    b2 = ax1.bar(xs+width/2, ev["recibidos"], width, label="Recibidos", color=PASTEL["verde"], alpha=0.9)
    b3 = ax1.bar(xs+width/2, ev["no_llego"], width, label="No Llego", color=PASTEL["naranja"], alpha=0.85, bottom=ev["recibidos"])
    b4 = ax1.bar(xs+width/2, ev["rechazados"], width, label="Rechazados", color=PASTEL["rojo"], alpha=0.85, bottom=ev["recibidos"]+ev["no_llego"])
    for rect in b1:
        h = rect.get_height()
        if h > 0: ax1.annotate(f"{h:.0f}", xy=(rect.get_x()+rect.get_width()/2, h), xytext=(0,3), textcoords="offset points", ha="center", va="bottom", fontsize=7, fontweight="bold", color="#374151")
    for rect in b2:
        h = rect.get_height()
        if h > 0: ax1.annotate(f"{h:.0f}", xy=(rect.get_x()+rect.get_width()/2, h), xytext=(0,3), textcoords="offset points", ha="center", va="bottom", fontsize=7, fontweight="bold", color="#16A34A")
    _ax_style(ax1)
    ax1.set_title("Agendados vs Recibidos por Dia", color="#1D4ED8", fontsize=11, fontweight="bold", pad=6)
    ax1.set_xticks(xs); ax1.set_xticklabels(lbl, rotation=45, ha="right", fontsize=7.5)
    ax1.legend(fontsize=8, labelcolor="#374151", facecolor="#F0F7FF", edgecolor="#DDE6F0", ncol=4, loc="upper right")
    ax2.plot(xs, pct_d, color="#1D4ED8", linewidth=2.5, marker="o", markersize=5, markerfacecolor="#60A5FA", zorder=5)
    ax2.fill_between(xs, pct_d, alpha=0.18, color="#60A5FA")
    ax2.axhline(85, color="#DC2626", linestyle="--", linewidth=1.5, label="Meta 85%", zorder=4)
    for x, v in zip(xs, pct_d):
        ax2.annotate(f"{v:.0f}%", xy=(x, v), xytext=(0, 7), textcoords="offset points", ha="center", fontsize=7.5, fontweight="bold", color="#1D4ED8")
    _ax_style(ax2)
    ax2.set_title("% Efectividad Diaria", color="#1D4ED8", fontsize=11, fontweight="bold", pad=6)
    ax2.set_ylim(0, 115); ax2.set_xticks(xs)
    ax2.set_xticklabels(lbl, rotation=45, ha="right", fontsize=7.5)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"{v:.0f}%"))
    ax2.legend(fontsize=8, labelcolor="#374151", facecolor="#F0F7FF", edgecolor="#DDE6F0")
    _img_place_stream(sl, _fig_to_stream(fig), 0.3, 0.95, 12.75, 6.25)
    _add_logo(sl); _add_nav(sl, prs, 2)

    # SLIDE 4: PROVEEDORES
    sl = _blank(prs); slides_list.append(sl)
    _rect(sl, 0, 0, 13.33, 7.5, OP_BG)
    _rect(sl, 0, 0, 0.22, 7.5, OP_AZ2)
    _rect(sl, 0.22, 0, 13.11, 0.62, OP_AZ)
    _txt(sl, " ANALISIS DE PROVEEDORES", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=OP_WH, align=PP_ALIGN.LEFT)
    gp_op = df_ag.groupby("_proveedor").agg(
        visitas=("_estado","count"), recibidos=("_estado",lambda x:(x=="RECIBIDO").sum()),
        rechazados=("_estado",lambda x:(x=="RECHAZADO").sum()), no_llego=("_estado",lambda x:x.str.contains("LLEG",na=False).sum()),
        estibas=("_estibas","sum")
    ).reset_index()
    gp_op["pct_ef"] = gp_op["recibidos"]/gp_op["visitas"].replace(0,1)
    gp_op["incumplidos"] = gp_op["rechazados"]+gp_op["no_llego"]
    top_v = gp_op.nlargest(10,"visitas")
    top_i = gp_op.nlargest(10,"incumplidos")
    top_e = gp_op[gp_op["visitas"]>=3].nlargest(10,"pct_ef")
    def shrt(s,n=20): return s[:n]+"..." if len(s)>n else s
    fig,(ax1,ax2,ax3) = plt.subplots(1,3, figsize=(13,4.8), facecolor="#F0F7FF")
    fig.subplots_adjust(wspace=0.5)
    lv = [shrt(p) for p in top_v["_proveedor"]]
    ys = np.arange(len(lv))
    ax1.barh(ys-0.2, top_v["visitas"], 0.35, color=PASTEL["azul"], alpha=0.92, label="Visitas")
    ax1.barh(ys+0.2, top_v["recibidos"], 0.35, color=PASTEL["verde"], alpha=0.92, label="Recibidos")
    for i,(a,b) in enumerate(zip(top_v["visitas"], top_v["recibidos"])):
        ax1.text(a+0.15, i-0.2, str(a), va="center", fontsize=7, color="#374151", fontweight="bold")
        ax1.text(b+0.15, i+0.2, str(b), va="center", fontsize=7, color="#16A34A", fontweight="bold")
    _ax_style(ax1); ax1.set_yticks(ys); ax1.set_yticklabels(lv, fontsize=7.5)
    ax1.invert_yaxis()
    ax1.set_title("Mayor Volumen (Visitas vs Recibidos)", color="#1D4ED8", fontsize=9.5, fontweight="bold")
    ax1.legend(fontsize=7.5, facecolor="#F0F7FF", edgecolor="#DDE6F0")
    li = [shrt(p) for p in top_i["_proveedor"]]
    yi_ax = np.arange(len(li))
    cols_i = [PASTEL["rojo"] if v>10 else PASTEL["naranja"] for v in top_i["incumplidos"]]
    ax2.barh(yi_ax, top_i["incumplidos"], color=cols_i, alpha=0.9)
    for i,v in enumerate(top_i["incumplidos"]):
        ax2.text(v+0.1, i, str(v), va="center", fontsize=7, color="#374151", fontweight="bold")
    _ax_style(ax2); ax2.set_yticks(yi_ax); ax2.set_yticklabels(li, fontsize=7.5)
    ax2.invert_yaxis()
    ax2.set_title("Mayores Incumplidores (No Llego + Rechazados)", color="#DC2626", fontsize=9.5, fontweight="bold")
    if len(top_e) > 0:
        le = [shrt(p) for p in top_e["_proveedor"]]
        ye_ax = np.arange(len(le))
        cols_e = [PASTEL["verde"] if v>=.8 else PASTEL["naranja"] if v>=.6 else PASTEL["rojo"] for v in top_e["pct_ef"]]
        ax3.barh(ye_ax, top_e["pct_ef"]*100, color=cols_e, alpha=0.9)
        ax3.axvline(85, color="#DC2626", linestyle="--", linewidth=1.2, label="Meta 85%")
        for i,v in enumerate(top_e["pct_ef"]*100):
            ax3.text(v+0.5, i, f"{v:.0f}%", va="center", fontsize=7, color="#374151", fontweight="bold")
        _ax_style(ax3); ax3.set_yticks(ye_ax); ax3.set_yticklabels(le, fontsize=7.5)
        ax3.invert_yaxis(); ax3.set_xlim(0, 110)
        ax3.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"{v:.0f}%"))
        ax3.legend(fontsize=7.5, facecolor="#F0F7FF", edgecolor="#DDE6F0")
        ax3.set_title("Mas Efectivos (>=3 visitas) % Efectividad", color="#1D4ED8", fontsize=9.5, fontweight="bold")
    _img_place_stream(sl, _fig_to_stream(fig), 0.25, 0.68, 12.85, 5.55)
    _add_logo(sl); _add_nav(sl, prs, 3)

    # SLIDE 5: NOVEDADES & HALLAZGOS
    sl = _blank(prs); slides_list.append(sl)
    _rect(sl, 0, 0, 13.33, 7.5, OP_BG)
    _rect(sl, 0, 0, 0.22, 7.5, OP_AZ2)
    _rect(sl, 0.22, 0, 13.11, 0.62, OP_AZ)
    _txt(sl, " NOVEDADES - TIEMPOS - HALLAZGOS", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=OP_WH, align=PP_ALIGN.LEFT)
    nov = df_cc.groupby("_novedad").size().sort_values(ascending=False).reset_index(name="n")
    nv_l = nov["_novedad"].str.strip().tolist()
    nv_v = nov["n"].tolist()
    nv_c = []
    for n in nv_l:
        if "SIN NOVEDAD" in n: nv_c.append(PASTEL["verde"])
        elif any(k in n for k in ["SIN PLACA","LLEGA TARDE","SIN HORARIO"]): nv_c.append(PASTEL["naranja"])
        else: nv_c.append(PASTEL["rojo"])
    df_t2 = df_cc[df_cc["_min"]>0].copy()
    df_t2["_ds"] = df_t2["_fecha"].astype(str)
    t_ev = df_t2.groupby("_ds")["_min"].mean().reset_index()
    fig,(ax1,ax2) = plt.subplots(1,2, figsize=(13,4.8), facecolor="#F0F7FF")
    ys_n = np.arange(len(nv_l))
    ax1.barh(ys_n, nv_v, color=nv_c, alpha=0.9)
    for i,(v,n) in enumerate(zip(nv_v,nv_l)):
        pct_s = f" ({v/len(df_cc)*100:.0f}%)" if len(df_cc) else ""
        ax1.text(v+0.15, i, f"{v}{pct_s}", va="center", fontsize=8, color="#374151", fontweight="bold")
    _ax_style(ax1); ax1.set_yticks(ys_n)
    ax1.set_yticklabels([n[:28] for n in nv_l], fontsize=8)
    ax1.invert_yaxis()
    ax1.set_title("Novedades de Descargue (CC)", color="#1D4ED8", fontsize=10.5, fontweight="bold")
    xs2 = np.arange(len(t_ev))
    vals_t = [round(v,1) for v in t_ev["_min"]]
    ax2.plot(xs2, vals_t, color="#EA580C", linewidth=2.5, marker="o", markersize=6, markerfacecolor="#FFCC80", zorder=5)
    ax2.fill_between(xs2, vals_t, alpha=0.18, color="#FFCC80")
    ax2.axhline(60, color="#DC2626", linestyle="--", linewidth=1.5, label="Limite 60 min")
    for x,v in zip(xs2, vals_t):
        ax2.annotate(f"{v:.0f}", xy=(x,v), xytext=(0,7), textcoords="offset points", ha="center", fontsize=7.5, fontweight="bold", color="#EA580C")
    _ax_style(ax2)
    ax2.set_title("Tiempo Prom. Descargue por Dia (min)", color="#1D4ED8", fontsize=10.5, fontweight="bold")
    ax2.set_xticks(xs2)
    ax2.set_xticklabels([d[5:] for d in t_ev["_ds"]], rotation=45, ha="right", fontsize=7.5)
    ax2.legend(fontsize=8, facecolor="#F0F7FF", edgecolor="#DDE6F0")
    _img_place_stream(sl, _fig_to_stream(fig), 0.25, 0.7, 12.85, 5.55)
    _add_logo(sl); _add_nav(sl, prs, 4)

    # GUARDAR
    fname_safe = re.sub(r'[/\\:*?"<>|]', '-', label).replace(' ','_')
    fname = f"PPTX_OPERATIVO_{fname_safe}_{date.today().strftime('%Y%m%d')}.pptx"
    ruta = os.path.join(carpeta, fname)
    prs.save(ruta)
    return ruta

# ═════════════════════════════════════════════════════════════════════════════
# PRESENTACIÓN GERENCIAL
# ═════════════════════════════════════════════════════════════════════════════
def generar_pptx_gerencial(df_ag, df_cc, label, fd, fh, carpeta):
    prs = _new_prs()
    total = len(df_ag)
    recib = int((df_ag["_estado"]=="RECIBIDO").sum())
    rechaz = int((df_ag["_estado"]=="RECHAZADO").sum())
    no_ll = int(df_ag["_estado"].str.contains("LLEG",na=False).sum())
    pct_ef = recib/total if total else 0
    est_tot = int(df_ag["_estibas"].sum())
    min_ok = df_cc[df_cc["_min"]>0]["_min"]
    prom_t = min_ok.mean() if len(min_ok) else 0
    sin_nov = int((df_cc["_novedad"].str.strip()=="SIN NOVEDAD").sum())
    pct_ok = sin_nov/len(df_cc) if len(df_cc) else 0
    fecha_g = date.today().strftime("%d/%m/%Y")
    sin_placa = int(df_cc["_novedad"].str.strip().eq("SIN PLACA").sum())
    no_msd = int(df_cc["_novedad"].str.strip().eq("NO DEJA AGENDAR MSD").sum())
    llega_t = int(df_cc["_novedad"].str.strip().eq("LLEGA TARDE").sum())

    # gp se calcula aqui para estar disponible en todas las diapositivas,
    # incluida la ultima (Conclusiones) que necesita sin_cc
    gp = df_ag.groupby("_proveedor").agg(
        visitas=("_estado","count"),
        recibidos=("_estado",lambda x:(x=="RECIBIDO").sum()),
        rechazados=("_estado",lambda x:(x=="RECHAZADO").sum()),
        no_llego=("_estado",lambda x:x.str.contains("LLEG",na=False).sum()),
        estibas=("_estibas","sum")
    ).reset_index()
    gp["pct_ef"] = gp["recibidos"]/gp["visitas"].replace(0,1)
    gp["incumplidos"] = gp["rechazados"]+gp["no_llego"]

    # Proveedores sin registro en control citas (se usa en slide 5)
    cc_prov = set(df_cc["_proveedor"].str.strip().str.upper())
    sin_cc = [p for p in gp["_proveedor"].str.upper() if p not in cc_prov]

    def _slide():
        sl = _blank(prs)
        _rect(sl, 0, 0, 13.33, 7.5, GR_BG)
        _rect(sl, 0, 0, 0.22, 7.5, GR_GD)
        return sl

    # SLIDE 1: PORTADA GERENCIAL
    sl = _slide()
    _rect(sl, 0, 0, 13.33, 3.4, GR_AZ)
    _img_place_file(sl, "assets/Imagen1.png", 0.5, 0.3, 1.85, 2.6)
    _txt(sl, "INFORME GERENCIAL", 2.6, 0.4, 10.3, 0.9, size=33, bold=True, color=GR_GD, align=PP_ALIGN.LEFT)
    _txt(sl, "RECEPCION DE MATERIALES - PLANTA ENTRERRIOS", 2.6, 1.32, 10.3, 0.48, size=15, color=GR_WH, align=PP_ALIGN.LEFT)
    _txt(sl, f"Periodo: {label} · {fd.strftime('%d/%m/%Y')} al {fh.strftime('%d/%m/%Y')}", 2.6, 1.85, 10.3, 0.38, size=11, color=RGBColor(0x90,0xB8,0xD8), align=PP_ALIGN.LEFT)
    _rect(sl, 0.35, 3.55, 12.63, 0.05, GR_GD)
    idx_items = ["1 - Portada","2 - KPIs Ejecutivos","3 - Tendencia & Evolucion","4 - Analisis Estrategico","5 - Conclusiones"]
    for i, item in enumerate(idx_items):
        yi = 3.75 + i * 0.58
        _rect(sl, 0.5, yi, 0.05, 0.38, GR_GD)
        _txt(sl, item, 0.72, yi+0.03, 5.5, 0.35, size=11, color=GR_AZ, align=PP_ALIGN.LEFT)
    _txt(sl, f"Generado: {fecha_g} | Agenda: {total} reg. | CC: {len(df_cc)} reg.", 0.35, 7.1, 12.63, 0.3, size=8.5, color=GR_GD, align=PP_ALIGN.CENTER, italic=True)
    _add_logo(sl); _add_nav(sl, prs, 0)

    # SLIDE 2: KPIs EJECUTIVOS
    sl = _slide()
    _rect(sl, 0.22, 0, 13.11, 0.62, GR_AZ)
    _txt(sl, " KPIs EJECUTIVOS - PERIODO", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=GR_GD, align=PP_ALIGN.LEFT)
    _txt(sl, f"{label} · {fd.strftime('%d/%m/%Y')} al {fh.strftime('%d/%m/%Y')}", 0.35, 0.65, 12, 0.3, size=9.5, color=GR_GD, align=PP_ALIGN.CENTER, italic=True)
    _img_place_file(sl, "assets/Gemini_Generated_Image_vcna55vcna55vcna.png", 7.8, 0.95, 5.4, 5.8)
    ef_col = GR_VD if pct_ef>=.85 else (OP_NR if pct_ef>=.70 else GR_RJ)
    ef_bg = RGBColor(0xD1,0xFA,0xE5) if pct_ef>=.85 else (RGBColor(0xFF,0xED,0xD5) if pct_ef>=.70 else RGBColor(0xFE,0xE2,0xE2))
    kpis_g = [
        (f"{pct_ef:.1%}", "Efectividad", ("OK Meta >=85%" if pct_ef>=.85 else "Bajo meta"), ef_bg, ef_col),
        (f"{prom_t:.0f} min", "Prom. Descargue", ("OK <=60 min" if prom_t<=60 else "Supera 60 min"),
         RGBColor(0xD1,0xFA,0xE5) if prom_t<=60 else RGBColor(0xFF,0xED,0xD5), GR_VD if prom_t<=60 else OP_NR),
        (f"{pct_ok:.1%}", "Sin Novedad CC", f"{sin_nov}/{len(df_cc)} reg.", GR_GDL, GR_GD),
        (str(no_msd), "Sin Agenda MSD", "Riesgo trazabilidad",
         RGBColor(0xFE,0xE2,0xE2) if no_msd>0 else RGBColor(0xD1,0xFA,0xE5), GR_RJ if no_msd>0 else GR_VD),
        (f"{est_tot:,}", "Estibas Recibidas", "Total periodo", GR_GDL, GR_GD),
    ]
    cols_per_row = 3
    for i,(v,l,s,bg,fc) in enumerate(kpis_g):
        row, col_i = divmod(i, cols_per_row)
        xi = 0.45 + col_i * 2.45
        yi = 1.05 + row * 1.68
        _rect(sl, xi, yi, 2.3, 1.55, bg)
        _rect(sl, xi, yi, 2.3, 0.08, fc)
        _txt(sl, v, xi+0.08, yi+0.15, 2.14, 0.72, size=27, bold=True, color=fc, align=PP_ALIGN.CENTER)
        _txt(sl, l, xi+0.06, yi+0.84, 2.18, 0.36, size=9, bold=True, color=RGBColor(0x37,0x41,0x51), align=PP_ALIGN.CENTER)
        if s:
            _txt(sl, s, xi+0.06, yi+1.17, 2.18, 0.3, size=7.5, color=GR_GD, align=PP_ALIGN.CENTER, italic=True)
    _add_logo(sl); _add_nav(sl, prs, 1)

    # SLIDE 3: EVOLUCIÓN Y TENDENCIA
    sl = _slide()
    _rect(sl, 0.22, 0, 13.11, 0.62, GR_AZ)
    _txt(sl, " TENDENCIA & EVOLUCION DEL PERIODO", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=GR_GD, align=PP_ALIGN.LEFT)
    df_ag2 = df_ag.copy()
    df_ag2["_mes"] = pd.to_datetime(df_ag2["_fecha"].astype(str)).dt.to_period("M")
    evo = df_ag2.groupby("_mes").agg(
        total=("_estado","count"), recibidos=("_estado", lambda x:(x=="RECIBIDO").sum()),
        no_llego=("_estado", lambda x:x.str.contains("LLEG",na=False).sum()),
        rechazados=("_estado",lambda x:(x=="RECHAZADO").sum()), estibas=("_estibas","sum")
    ).reset_index()
    ml = [str(m) for m in evo["_mes"]]
    xs = np.arange(len(ml))
    mp = [r/t*100 if t else 0 for r,t in zip(evo["recibidos"],evo["total"])]
    fig,(ax1,ax2) = plt.subplots(1,2, figsize=(12.5,5.1), facecolor="#FFF8F0")
    fig.subplots_adjust(wspace=0.38)
    w = 0.38
    b1 = ax1.bar(xs-w/2, evo["total"], w, label="Agendados", color=PASTEL["celeste"], alpha=0.9)
    b2 = ax1.bar(xs+w/2, evo["recibidos"], w, label="Recibidos", color=PASTEL["verde"], alpha=0.9)
    b3 = ax1.bar(xs+w/2, evo["no_llego"], w, label="No Llego", color=PASTEL["naranja"], alpha=0.85, bottom=evo["recibidos"])
    for rect in b1:
        h = rect.get_height()
        if h > 0: ax1.annotate(f"{h:.0f}", xy=(rect.get_x()+rect.get_width()/2, h), xytext=(0,3), textcoords="offset points", ha="center", va="bottom", fontsize=8, fontweight="bold", color="#374151")
    for rect in b2:
        h = rect.get_height()
        if h > 0: ax1.annotate(f"{h:.0f}", xy=(rect.get_x()+rect.get_width()/2, h), xytext=(0,3), textcoords="offset points", ha="center", va="bottom", fontsize=8, fontweight="bold", color="#16A34A")
    _ax_style_warm(ax1)
    ax1.set_xticks(xs); ax1.set_xticklabels(ml, rotation=0, fontsize=9.5)
    ax1.set_title("Volumen Mensual", color="#1E3A8A", fontsize=11, fontweight="bold")
    ax1.legend(fontsize=8, facecolor="#FFF8F0", edgecolor="#F0E0D0")
    cols_m = [PASTEL["verde"] if v>=85 else PASTEL["naranja"] if v>=70 else PASTEL["rojo"] for v in mp]
    bars = ax2.bar(xs, mp, color=cols_m, alpha=0.9, width=0.6, edgecolor="#FFFFFF", linewidth=0.8)
    ax2.axhline(85, color="#DC2626", linestyle="--", linewidth=1.8, label="Meta 85%")
    for x,v in zip(xs,mp):
        ax2.annotate(f"{v:.0f}%", xy=(x,v), xytext=(0,5), textcoords="offset points", ha="center", fontsize=9.5, fontweight="bold", color="#1E3A8A")
    _ax_style_warm(ax2)
    ax2.set_xticks(xs); ax2.set_xticklabels(ml, rotation=0, fontsize=9.5)
    ax2.set_title("% Efectividad Mensual", color="#1E3A8A", fontsize=11, fontweight="bold")
    ax2.set_ylim(0, 115)
    ax2.legend(fontsize=8, facecolor="#FFF8F0", edgecolor="#F0E0D0")
    _img_place_stream(sl, _fig_to_stream(fig), 0.3, 0.7, 12.75, 6.1)
    _add_logo(sl); _add_nav(sl, prs, 2)

    # SLIDE 4: PROVEEDORES ESTRATÉGICOS
    sl = _slide()
    _rect(sl, 0.22, 0, 13.11, 0.62, GR_AZ)
    _txt(sl, " ANALISIS ESTRATEGICO DE PROVEEDORES", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=GR_GD, align=PP_ALIGN.LEFT)
    top_e = gp[gp["visitas"]>=3].nlargest(10,"pct_ef")
    top_i = gp.nlargest(10,"incumplidos")
    def shn(s,n=22): return s[:n]+"..." if len(s)>n else s
    fig,(ax1,ax2) = plt.subplots(1,2, figsize=(12.5,5.1), facecolor="#FFF8F0")
    fig.subplots_adjust(wspace=0.45)
    if len(top_e) > 0:
        le = [shn(p) for p in top_e["_proveedor"]]
        ye = np.arange(len(le))
        cols_e = [PASTEL["verde"] if v>=.8 else PASTEL["naranja"] if v>=.6 else PASTEL["rojo"] for v in top_e["pct_ef"]]
        ax1.barh(ye, top_e["pct_ef"]*100, color=cols_e, alpha=0.9, edgecolor="#FFFFFF")
        ax1.axvline(85, color="#DC2626", linestyle="--", linewidth=1.2, label="Meta 85%")
        for i,v in enumerate(top_e["pct_ef"]*100):
            ax1.text(v+0.5, i, f"{v:.0f}%", va="center", fontsize=8, color="#374151", fontweight="bold")
        _ax_style_warm(ax1)
        ax1.set_yticks(ye); ax1.set_yticklabels(le, fontsize=8)
        ax1.invert_yaxis(); ax1.set_xlim(0,112)
        ax1.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_:f"{v:.0f}%"))
        ax1.legend(fontsize=7.5, facecolor="#FFF8F0", edgecolor="#F0E0D0")
        ax1.set_title("Proveedores mas Efectivos (>=3 visitas)", color="#1E3A8A", fontsize=10, fontweight="bold")
    li = [shn(p) for p in top_i["_proveedor"]]
    yi2 = np.arange(len(li))
    cols_i = [PASTEL["rojo"] if v>10 else PASTEL["naranja"] for v in top_i["incumplidos"]]
    ax2.barh(yi2, top_i["incumplidos"], color=cols_i, alpha=0.9, edgecolor="#FFFFFF")
    for i,v in enumerate(top_i["incumplidos"]):
        ax2.text(v+0.1, i, str(v), va="center", fontsize=8, color="#374151", fontweight="bold")
    _ax_style_warm(ax2)
    ax2.set_yticks(yi2); ax2.set_yticklabels(li, fontsize=8)
    ax2.invert_yaxis()
    ax2.set_title("Mayores Incumplidores", color="#DC2626", fontsize=10, fontweight="bold")
    _img_place_stream(sl, _fig_to_stream(fig), 0.3, 0.7, 12.75, 6.1)
    _add_logo(sl); _add_nav(sl, prs, 3)

    # SLIDE 5: CONCLUSIONES ESTRATÉGICAS
    sl = _slide()
    _rect(sl, 0.22, 0, 13.11, 0.62, GR_AZ)
    _txt(sl, " CONCLUSIONES & RECOMENDACIONES", 0.35, 0.08, 12, 0.46, size=18, bold=True, color=GR_GD, align=PP_ALIGN.LEFT)

    concl = [
        (pct_ef>=.85,
         f"Efectividad: {pct_ef:.1%}",
         "Cumple meta >=85%." if pct_ef>=.85 else f"Bajo meta - analizar No Llego ({no_ll} veh.) como causa principal."),
        (no_msd==0,
         f"Agenda MSD: {no_msd} sin registro",
         "Sin riesgo de trazabilidad." if no_msd==0 else "Implementar control MSD antes de autorizar ingreso."),
        (sin_placa==0,
         f"Sin placa: {sin_placa} vehiculos",
         "Documentacion completa." if sin_placa==0 else "Exigir placa como requisito obligatorio de agendamiento."),
        (prom_t<=60,
         f"T. prom. descargue: {prom_t:.0f} min",
         "Dentro del objetivo <=60 min." if prom_t<=60 else "Revisar recursos de descargue y balanceo de muelles."),
        (llega_t==0,
         f"Llegan tarde: {llega_t} vehiculos",
         "Puntualidad adecuada." if llega_t==0 else "Implementar ventanas de tiempo mas estrictas."),
        (len(sin_cc)==0,
         f"Sin registro CC: {len(sin_cc)} proveedores",
         "Cobertura completa." if not sin_cc else f"Registrar {len(sin_cc)} proveedor(es) en Control de Citas."),
    ]

    for i, (ok, titulo, detalle) in enumerate(concl):
        row, ci = divmod(i, 2)
        xi = 0.38 if ci == 0 else 6.9
        yi = 0.78 + row * 2.1
        bg_c = RGBColor(0xF0,0xFD,0xF4) if ok else RGBColor(0xFF,0xF1,0xF1)
        bd_c = GR_VD if ok else GR_RJ
        _rect(sl, xi, yi, 6.1, 1.85, bg_c)
        _rect(sl, xi, yi, 0.11, 1.85, bd_c)
        # Usar simbolos BMP para evitar errores de codificacion XML en python-pptx
        ico = "OK" if ok else "!!"
        _txt(sl, f"[{ico}] {titulo}", xi+0.2, yi+0.12, 5.8, 0.48,
             size=11, bold=True, color=bd_c)
        _txt(sl, detalle, xi+0.2, yi+0.58, 5.8, 0.92,
             size=9.5, color=RGBColor(0x37,0x41,0x51))

    _txt(sl, f"Generado: {fecha_g} | {label}", 0.35, 7.1, 12.63, 0.3,
         size=8, color=GR_GD, align=PP_ALIGN.CENTER, italic=True)
    _add_logo(sl)
    _add_nav(sl, prs, 4)

    # GUARDAR
    fname_safe = re.sub(r'[/\\:*?"<>|]', '-', label).replace(' ','_')
    fname = f"PPTX_GERENCIAL_{fname_safe}_{date.today().strftime('%Y%m%d')}.pptx"
    ruta = os.path.join(carpeta, fname)
    prs.save(ruta)
    return ruta
